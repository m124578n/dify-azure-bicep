"""Generate deployment proposal PPTX for management presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color Palette ──────────────────────────────────────────────────────────────
AZURE_BLUE   = RGBColor(0x00, 0x78, 0xD4)
DARK_GRAY    = RGBColor(0x23, 0x23, 0x23)
MID_GRAY     = RGBColor(0x60, 0x60, 0x60)
LIGHT_GRAY   = RGBColor(0xF3, 0xF3, 0xF3)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GREEN        = RGBColor(0x10, 0x79, 0x47)
ORANGE       = RGBColor(0xD8, 0x34, 0x00)
ACCENT_BLUE  = RGBColor(0x00, 0x5A, 0x9E)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # blank layout


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def txbox(slide, text, l, t, w, h, size=18, bold=False, color=DARK_GRAY,
          align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Microsoft JhengHei"
    return box


def header_bar(slide, title, subtitle=None):
    """Blue header bar at top."""
    rect(slide, 0, 0, 13.33, 1.4, fill=AZURE_BLUE)
    txbox(slide, title, 0.4, 0.15, 10, 0.7, size=28, bold=True, color=WHITE)
    if subtitle:
        txbox(slide, subtitle, 0.4, 0.82, 10, 0.45, size=14, color=RGBColor(0xCC, 0xE4, 0xF7))


def footer(slide, page_num, total=9):
    rect(slide, 0, 7.15, 13.33, 0.35, fill=RGBColor(0xE8, 0xF1, 0xFB))
    txbox(slide, "Dify on Azure — 部署架構說明", 0.3, 7.17, 9, 0.28,
          size=9, color=MID_GRAY)
    txbox(slide, f"{page_num} / {total}", 12.3, 7.17, 0.8, 0.28,
          size=9, color=MID_GRAY, align=PP_ALIGN.RIGHT)


def bullet_box(slide, items, l, t, w, h, title=None, title_color=AZURE_BLUE,
               bullet="●", indent=0, size=14):
    """Multi-line bullet list inside a box."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(size + 1)
        run.font.bold = True
        run.font.color.rgb = title_color
        run.font.name = "Microsoft JhengHei"
    for item in items:
        p = tf.add_paragraph() if not first else (tf.paragraphs[0] if not title else tf.add_paragraph())
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        run = p.add_run()
        prefix = "  " * indent + bullet + "  " if bullet else "  " * indent
        run.text = prefix + item
        run.font.size = Pt(size)
        run.font.color.rgb = DARK_GRAY
        run.font.name = "Microsoft JhengHei"


def card(slide, l, t, w, h, title, items, icon=None,
         title_color=AZURE_BLUE, bg=LIGHT_GRAY):
    rect(slide, l, t, w, h, fill=bg, line=RGBColor(0xDD, 0xDD, 0xDD))
    head = (icon + "  " if icon else "") + title
    txbox(slide, head, l + 0.15, t + 0.1, w - 0.3, 0.38,
          size=13, bold=True, color=title_color)
    # divider
    d = slide.shapes.add_shape(1,
        Inches(l + 0.1), Inches(t + 0.48),
        Inches(w - 0.2), Inches(0.02))
    d.fill.solid(); d.fill.fore_color.rgb = title_color
    d.line.fill.background()
    y = t + 0.55
    for item in items:
        txbox(slide, "·  " + item, l + 0.15, y, w - 0.3, 0.3,
              size=11, color=DARK_GRAY)
        y += 0.28


# ══════════════════════════════════════════════════════════════════════════════
# Slide 1 — Cover
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=AZURE_BLUE)
rect(s, 0, 4.8, 13.33, 2.7, fill=ACCENT_BLUE)

txbox(s, "Dify on Azure", 1, 1.2, 11, 1.1, size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, "部署架構說明  ·  Bicep IaC 方案", 1, 2.4, 11, 0.7,
      size=22, color=RGBColor(0xCC, 0xE4, 0xF7), align=PP_ALIGN.CENTER)
txbox(s, "Infrastructure as Code  |  Azure Container Apps  |  CI/CD 自動化",
      1, 3.1, 11, 0.5, size=14, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)
txbox(s, "SI Team  ·  2026", 1, 5.2, 11, 0.5,
      size=13, color=RGBColor(0xCC, 0xE4, 0xF7), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# Slide 2 — Agenda
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "本次說明重點")
footer(s, 2)

items = [
    ("01", "什麼是 Bicep？",          "Azure 官方 IaC 工具介紹"),
    ("02", "為什麼選擇 Bicep？",       "相較手動部署與其他工具的優勢"),
    ("03", "Dify 部署架構",            "Azure 服務組成與網路設計"),
    ("04", "deploy.ps1 部署流程",      "一鍵完成基礎架構部署與初始化"),
    ("05", "所需權限說明",             "最小權限原則與安全疑慮釐清"),
    ("06", "安全控管機制",             "程式碼審計、版控與密碼管理"),
    ("07", "CI/CD 自動化流程",         "Azure DevOps Pipelines 自動部署"),
    ("08", "風險與緩解措施",           "常見疑慮與對應做法"),
]
for i, (num, title, desc) in enumerate(items):
    row = 1.6 + i * 0.7
    rect(s, 0.5, row, 0.55, 0.52, fill=AZURE_BLUE)
    txbox(s, num, 0.5, row + 0.06, 0.55, 0.4, size=14, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s, title, 1.2, row + 0.04, 3.5, 0.28, size=14, bold=True, color=DARK_GRAY)
    txbox(s, desc,  1.2, row + 0.28, 5,   0.25, size=11, color=MID_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 3 — 什麼是 Bicep
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "什麼是 Bicep？", "Azure 官方 Infrastructure as Code 工具")
footer(s, 3)

txbox(s, "Infrastructure as Code（IaC）是將雲端基礎架構的建置、設定以程式碼描述的方式管理，"
         "取代人工在 Portal 逐步點選操作。",
      0.5, 1.55, 12.3, 0.55, size=13, color=MID_GRAY)

cards_data = [
    ("🔷 Azure 原生", ["微軟官方開發與維護", "直接對應 Azure Resource Manager", "完整 Azure 服務支援"]),
    ("📝 宣告式語法", ["描述「要什麼」而非「怎麼做」", "語法比 ARM Template 簡潔 60%+", "VS Code 有完整 IntelliSense"]),
    ("♻️ 冪等部署", ["多次執行結果一致", "只更新有差異的資源", "不會重複建立已存在的資源"]),
    ("🔍 可審計", ["所有設定都在程式碼中", "Git 版控，每次變更有紀錄", "Code Review 流程可控管"]),
]
positions = [(0.5, 2.3), (3.6, 2.3), (6.7, 2.3), (9.8, 2.3)]
for (l, t), (title, items) in zip(positions, cards_data):
    card(s, l, t, 2.9, 4.5, title, items)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 4 — 為什麼選擇 Bicep
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "為什麼選擇 Bicep？", "與其他方式的比較")
footer(s, 4)

# Table header
rect(s, 0.4, 1.55, 12.5, 0.45, fill=AZURE_BLUE)
for x, label in [(0.5, "比較項目"), (3.2, "手動（Portal）"), (5.9, "PowerShell Script"),
                 (8.6, "Terraform"), (11.3, "Bicep ✓")]:
    txbox(s, label, x, 1.6, 2.5, 0.35, size=12, bold=True, color=WHITE)

rows = [
    ("可重複部署",      "✗ 需手動重做",  "△ 需維護腳本",  "✓",          "✓"),
    ("Git 版本控制",    "✗",             "✓",              "✓",          "✓"),
    ("Azure 原生支援",  "✓",             "△",              "△ 需 Provider", "✓ 官方"),
    ("State 管理",      "─",             "─",              "✗ 需 backend",  "✓ 內建"),
    ("學習成本",        "低",            "中",             "中高",        "低（類 JSON）"),
    ("適合 Azure 專案", "△",             "△",              "✓",          "✓ 最佳"),
]
row_colors = [WHITE, LIGHT_GRAY]
for i, (item, *cols) in enumerate(rows):
    y = 2.05 + i * 0.65
    rect(s, 0.4, y, 12.5, 0.62, fill=row_colors[i % 2])
    txbox(s, item, 0.5, y + 0.15, 2.5, 0.35, size=12, bold=True, color=DARK_GRAY)
    for j, val in enumerate(cols):
        color = GREEN if val.startswith("✓") else (ORANGE if val.startswith("✗") else DARK_GRAY)
        txbox(s, val, 3.2 + j * 2.7, y + 0.15, 2.5, 0.35, size=12, color=color)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 5 — 部署架構
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "Dify 部署架構", "Azure Container Apps + 周邊服務")
footer(s, 5)

# Resource Group box
rect(s, 0.3, 1.55, 12.7, 5.4, fill=RGBColor(0xF0, 0xF7, 0xFF),
     line=AZURE_BLUE)
txbox(s, "Resource Group：rg-Dify", 0.5, 1.6, 5, 0.35, size=11,
      bold=True, color=AZURE_BLUE)

# VNet box
rect(s, 0.5, 2.05, 9.0, 4.6, fill=RGBColor(0xE8, 0xF5, 0xE9),
     line=GREEN)
txbox(s, "Virtual Network（VNet）", 0.65, 2.1, 4, 0.3, size=10,
      bold=True, color=GREEN)

# ACA Environment
rect(s, 0.7, 2.5, 5.8, 3.8, fill=WHITE, line=AZURE_BLUE)
txbox(s, "Azure Container Apps Environment", 0.85, 2.55, 5, 0.3,
      size=10, bold=True, color=AZURE_BLUE)

aca_apps = [
    ("nginx",   "對外入口", 0.85, 2.95),
    ("api",     "API Server", 0.85, 3.35),
    ("web",     "前端 UI", 0.85, 3.75),
    ("worker",  "背景任務", 3.3, 2.95),
    ("sandbox", "程式執行", 3.3, 3.35),
    ("plugin",  "Plugin Daemon", 3.3, 3.75),
]
for name, desc, x, y in aca_apps:
    rect(s, x, y, 2.2, 0.35, fill=RGBColor(0xE3, 0xF2, 0xFD),
         line=AZURE_BLUE)
    txbox(s, f"{name}  |  {desc}", x + 0.08, y + 0.05, 2.1, 0.28,
          size=9.5, color=DARK_GRAY)

# Postgres + Redis in VNet
rect(s, 0.7, 4.2, 5.8, 1.7, fill=RGBColor(0xFF, 0xF8, 0xE1),
     line=RGBColor(0xF5, 0xA6, 0x23))
txbox(s, "Private Subnet", 0.85, 4.25, 3, 0.28, size=9,
      bold=True, color=RGBColor(0xC0, 0x70, 0x00))
rect(s, 0.85, 4.6, 2.5, 0.55, fill=WHITE, line=RGBColor(0xF5, 0xA6, 0x23))
txbox(s, "PostgreSQL Flexible Server\n+ pgvector", 0.95, 4.65, 2.3, 0.45, size=9, color=DARK_GRAY)
rect(s, 3.5, 4.6, 2.5, 0.55, fill=WHITE, line=RGBColor(0xF5, 0xA6, 0x23))
txbox(s, "Azure Cache for Redis", 3.6, 4.72, 2.3, 0.3, size=9, color=DARK_GRAY)

# Right side services
right_items = [
    ("Azure Storage\nAccount", 2.05),
    ("Azure Container\nRegistry（ACR）", 2.75),
    ("Log Analytics\nWorkspace", 3.45),
    ("Private DNS\nZone", 4.15),
]
for label, y in right_items:
    rect(s, 9.7, y, 3.1, 0.6, fill=LIGHT_GRAY, line=RGBColor(0xCC, 0xCC, 0xCC))
    txbox(s, label, 9.85, y + 0.05, 2.8, 0.5, size=10, color=DARK_GRAY)

txbox(s, "周邊服務", 9.7, 1.65, 3, 0.3, size=10, bold=True, color=MID_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 6 — deploy.ps1 部署流程
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "deploy.ps1 部署流程", "一鍵完成 Bicep 部署 + 設定檔上傳 + DB 初始化")
footer(s, 6)

_phases = [
    (AZURE_BLUE, "Phase 1\n基礎架構", [
        ("1", "Azure CLI 登入確認",  "az account show\naz login（如未登入）"),
        ("2", "建立 Resource Group", "az group create\n（已存在則略過）"),
        ("3", "Bicep 模板部署",      "az deployment group create\n--template-file main.bicep"),
    ]),
    (ORANGE, "Phase 2\n設定檔上傳", [
        ("4", "取得 Storage 金鑰",   "az storage account keys list"),
        ("5", "產生 SAS Token",      "有效期 24hr，https-only"),
        ("6", "上傳至 File Share",   "azcopy 優先 → az CLI fallback\nnginx / ssrfproxy / sandbox"),
    ]),
    (GREEN, "Phase 3\n服務啟動", [
        ("7", "還原儲存體安全設定",  "default-action Deny\nbypass AzureServices"),
        ("8", "重啟 nginx",          "containerapp revision restart"),
        ("9", "DB Migration",        "exec api → flask db upgrade"),
    ]),
]

_row_h = 1.5
_row_gap = 0.08
_step_w = 3.6
_step_gap = 0.07
_step_x0 = 2.05

for _r, (_ph_color, _ph_label, _ph_steps) in enumerate(_phases):
    _t = 1.6 + _r * (_row_h + _row_gap)
    rect(s, 0.3, _t, 1.65, _row_h, fill=_ph_color)
    txbox(s, _ph_label, 0.3, _t + _row_h / 2 - 0.3, 1.65, 0.7,
          size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for _c, (_num, _title, _desc) in enumerate(_ph_steps):
        _l = _step_x0 + _c * (_step_w + _step_gap)
        rect(s, _l, _t, _step_w, _row_h, fill=WHITE, line=RGBColor(0xCC, 0xCC, 0xCC))
        rect(s, _l, _t, 0.07, _row_h, fill=_ph_color)
        rect(s, _l + 0.12, _t + 0.1, 0.36, 0.36, fill=_ph_color)
        txbox(s, _num, _l + 0.12, _t + 0.1, 0.36, 0.36,
              size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txbox(s, _title, _l + 0.58, _t + 0.08, _step_w - 0.65, 0.5,
              size=11, bold=True, color=DARK_GRAY)
        txbox(s, _desc, _l + 0.12, _t + 0.65, _step_w - 0.2, 0.75,
              size=9.5, color=MID_GRAY)

txbox(s, "✓  完成：顯示 nginx Endpoint URL，Dify 服務上線",
      0.3, 6.35, 12.7, 0.35, size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 7 — 所需權限說明
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "所需權限說明", "為什麼需要 Subscription Contributor？")
footer(s, 7)

# One big answer
rect(s, 0.4, 1.55, 12.5, 0.55, fill=AZURE_BLUE)
txbox(s, "申請角色：Subscription Contributor　　（僅此一項）",
      0.6, 1.62, 11, 0.38, size=15, bold=True, color=WHITE)

# Root cause explanation
rect(s, 0.4, 2.2, 12.5, 1.5, fill=RGBColor(0xF0, 0xF7, 0xFF), line=AZURE_BLUE)
txbox(s, "根本原因", 0.6, 2.27, 3, 0.3, size=12, bold=True, color=AZURE_BLUE)
txbox(s,
      "Azure Container Apps 環境在啟用 VNet 整合時，Azure 平台會在背景自動建立"
      "一個 Infrastructure Resource Group（ACA Infra RG），並在其中部署內部網路元件"
      "（Load Balancer、NIC、Public IP 等）。\n\n"
      "這個 RG 是由 Azure 自行建立與管理，部署者無需手動操作，"
      "但 Azure 在執行此動作時會向 Subscription 層級要求 Contributor 授權。"
      "若缺少此權限，ACA 環境建立流程會卡住並回傳錯誤。",
      0.6, 2.6, 12.1, 1.0, size=12, color=DARK_GRAY)

# Three point cards
points = [
    ("🔍 影響範圍", AZURE_BLUE, [
        "權限作用於指定 Subscription",
        "實際資源只建立在 rg-Dify-SITeam",
        "ACA Infra RG 由 Azure 自動管理，部署者不直接操作",
    ]),
    ("🛡 安全疑慮釐清", RGBColor(0x10, 0x79, 0x47), [
        "Contributor 不能修改 RBAC（無法升權）",
        "不能存取其他既有 RG 的資源內容",
        "所有操作皆有 Azure Activity Log 可稽核",
    ]),
    ("📋 替代方案評估", ORANGE, [
        "RG-level Contributor：ACA 建立會失敗，不可行",
        "Custom Role：可限縮動作，但需額外設定",
        "Subscription Contributor：最直接，符合 Azure 建議做法",
    ]),
]
for i, (title, color, items) in enumerate(points):
    l = 0.4 + i * 4.3
    card(s, l, 3.85, 4.1, 3.15, title, items, title_color=color, bg=WHITE)

txbox(s, "⚠  Subscription Contributor 是 Azure Container Apps VNet 整合的平台需求，非過度授權。",
      0.5, 6.88, 12.3, 0.35, size=11, color=ORANGE)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 7 — 安全控管機制
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "安全控管機制", "程式碼審計、版控、密碼管理")
footer(s, 8)

sections = [
    ("🔒 密碼與 Secret 管理", GREEN, [
        "所有密碼參數標記 @secure()，部署時不會寫入 git",
        "parameters.json 已加入 .gitignore，只有本地存在",
        "正式環境建議串接 Azure Key Vault（可後續擴充）",
    ]),
    ("📋 程式碼審計", AZURE_BLUE, [
        "所有 Bicep 程式碼在 git，每次變更都有 commit 記錄",
        "部署前可執行 bicep build 靜態分析，確保語法正確",
        "Azure Repos PR / Code Review 可審核每次基礎架構變更",
    ]),
    ("🛡 網路安全", RGBColor(0x7A, 0x1C, 0xA0), [
        "所有服務部署於 VNet 內，對外只開放 nginx（port 443）",
        "PostgreSQL、Redis 透過 Private Endpoint，不對外暴露",
        "Storage Account 部署後關閉公開存取",
    ]),
    ("🔄 變更管理", ORANGE, [
        "基礎架構變更需經 Azure Repos PR 審核",
        "deploy.ps1 支援 dev / prod 環境分離",
        "Bicep 冪等部署，可隨時重新執行不影響現有資料",
    ]),
]
for i, (title, color, items) in enumerate(sections):
    col = i % 2
    row = i // 2
    l = 0.4 + col * 6.45
    t = 1.6 + row * 2.65
    card(s, l, t, 6.2, 2.45, title, items, title_color=color, bg=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 8 — CI/CD 流程
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "CI/CD 自動化流程", "Azure DevOps Pipelines + Azure Repos")
footer(s, 9)

# Two flow columns
for col, (title, color, steps) in enumerate([
    ("Plugin Image 版控", AZURE_BLUE, [
        "Plugin 程式碼修改",
        "本地 Unit Test",
        "git push → Azure Repos PR",
        "Docker Build & Push to ACR",
        "更新 ACA container image tag",
        "Dev 環境 Smoke Test",
        "Prod 部署",
    ]),
    ("Dify Workflow 版控", GREEN, [
        "開發人員在 UAT 修改 Workflow",
        "export --app-id → YAML 儲存至 git",
        "git push → PR Review",
        "merge to main",
        "Azure Pipelines 觸發",
        "deploy script → Prod 更新 draft",
        "自動 publish → Prod 上線",
    ]),
]):
    l = 0.4 + col * 6.45
    rect(s, l, 1.55, 6.2, 0.45, fill=color)
    txbox(s, title, l + 0.15, 1.6, 5.8, 0.35, size=14, bold=True, color=WHITE)
    for i, step in enumerate(steps):
        y = 2.1 + i * 0.68
        rect(s, l + 0.15, y, 0.45, 0.45, fill=color)
        txbox(s, str(i + 1), l + 0.15, y + 0.07, 0.45, 0.3,
              size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(s, l + 0.65, y + 0.05, 5.3, 0.38,
             fill=LIGHT_GRAY if i % 2 == 0 else WHITE,
             line=RGBColor(0xDD, 0xDD, 0xDD))
        txbox(s, step, l + 0.75, y + 0.1, 5.1, 0.3, size=11.5, color=DARK_GRAY)
        if i < len(steps) - 1:
            arr = slide.shapes if (slide := s) else None
            txbox(s, "↓", l + 0.27, y + 0.45, 0.25, 0.25,
                  size=10, color=color, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 9 — 風險與緩解
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "常見疑慮與對應做法", "風險評估與緩解措施")
footer(s, 9)

concerns = [
    (
        "Bicep 是新工具，團隊不熟悉",
        "學習曲線低，語法類似 JSON/YAML\n微軟官方文件完整，有 VS Code 擴充套件\n此次已完成基礎架構，維護成本低",
    ),
    (
        "Owner 權限風險過大",
        "Owner 僅限 ACA Infra RG（無業務資料）\n主要 RG 只需 Contributor\n可搭配 Azure Policy 限制可用資源種類",
    ),
    (
        "部署出錯難以回復",
        "Bicep 冪等，可重新執行恢復正確狀態\n所有變更有 git 記錄，可 rollback\n資料庫與儲存體不在 Bicep 刪除範圍內",
    ),
    (
        "密碼洩漏風險",
        "parameters.json 已 gitignore\n@secure() 參數不會出現在部署日誌\n後續可串接 Azure Key Vault 強化",
    ),
]

rect(s, 0.3, 1.5, 12.7, 0.38, fill=AZURE_BLUE)
for x, label in [(0.5, "疑慮"), (6.5, "對應做法")]:
    txbox(s, label, x, 1.55, 5.5, 0.28, size=12, bold=True, color=WHITE)

for i, (concern, mitigation) in enumerate(concerns):
    y = 2.0 + i * 1.2
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    rect(s, 0.3, y, 12.7, 1.12, fill=bg, line=RGBColor(0xDD, 0xDD, 0xDD))
    rect(s, 0.3, y, 0.25, 1.12, fill=ORANGE)
    txbox(s, "Q", 0.3, y + 0.35, 0.25, 0.4, size=13, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s, concern,    0.65, y + 0.1,  5.5, 0.9,  size=12, bold=True,  color=DARK_GRAY)
    txbox(s, mitigation, 6.5,  y + 0.08, 6.3, 0.95, size=11, color=MID_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
out = r"C:\Users\m23568n\projects\bicep\dify-azure-bicep\Dify_Azure_部署架構說明.pptx"
prs.save(out)
print(f"Saved: {out}")
